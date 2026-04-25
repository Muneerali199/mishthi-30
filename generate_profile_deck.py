from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.replace("#", "")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


BG = rgb("#F8FAFC")
WHITE = rgb("#FFFFFF")
DARK = rgb("#1F2937")
SLATE = rgb("#475569")
MUTED = rgb("#E2E8F0")
TEAL = rgb("#0F766E")
TEAL_SOFT = rgb("#ECFDF5")
AMBER = rgb("#F59E0B")
AMBER_SOFT = rgb("#FFFBEB")


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

prs.core_properties.title = "Mishthi | Python Developer Profile"
prs.core_properties.author = "OpenAI Codex"
prs.core_properties.subject = "Recruiter-friendly personal introduction deck"


def add_bg(slide, color=BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_blob(slide, x, y, w, h, color, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.fill.background()
    return shape


def add_card(slide, x, y, w, h, fill, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=18,
    color=DARK,
    font="Aptos",
    bold=False,
    align=PP_ALIGN.LEFT,
    italic=False,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_paragraphs(
    slide,
    lines,
    x,
    y,
    w,
    h,
    size=16,
    color=SLATE,
    font="Aptos",
    bullet=False,
    line_gap=6,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {line}" if bullet else line
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(line_gap)
        p.line_spacing = 1.15
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_label(slide, text, x, y, w=2.0, h=0.4):
    chip = add_card(slide, Inches(x), Inches(y), Inches(w), Inches(h), TEAL_SOFT)
    add_text(
        slide,
        text.upper(),
        Inches(x + 0.16),
        Inches(y + 0.05),
        Inches(w - 0.3),
        Inches(h - 0.1),
        size=10,
        color=TEAL,
        font="Aptos",
        bold=True,
    )
    return chip


def add_chip(slide, text, x, y, w, fill, text_color=WHITE):
    add_card(slide, Inches(x), Inches(y), Inches(w), Inches(0.42), fill)
    add_text(
        slide,
        text,
        Inches(x + 0.12),
        Inches(y + 0.05),
        Inches(w - 0.24),
        Inches(0.28),
        size=11,
        color=text_color,
        font="Aptos",
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_section_title(slide, kicker, title, subtitle=None, x=0.8, y=0.55, title_color=DARK):
    add_label(slide, kicker, x, y, w=1.75, h=0.38)
    add_text(
        slide,
        title,
        Inches(x),
        Inches(y + 0.48),
        Inches(6.4),
        Inches(0.9),
        size=27,
        color=title_color,
        font="Aptos Display",
        bold=True,
    )
    if subtitle:
        add_text(
            slide,
            subtitle,
            Inches(x),
            Inches(y + 1.28),
            Inches(6.8),
            Inches(0.7),
            size=13,
            color=SLATE if title_color == DARK else MUTED,
            font="Aptos",
        )


def add_divider(slide, x, y, w, color=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def slide_cover():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_blob(slide, Inches(-1.0), Inches(-0.7), Inches(2.8), Inches(2.8), TEAL_SOFT)
    add_blob(slide, Inches(11.6), Inches(5.9), Inches(2.3), Inches(2.3), AMBER_SOFT)
    add_blob(slide, Inches(10.2), Inches(-0.5), Inches(2.0), Inches(2.0), MUTED, transparency=0.2)

    add_text(
        slide,
        "Mishthi",
        Inches(0.85),
        Inches(1.0),
        Inches(5.5),
        Inches(0.9),
        size=31,
        color=DARK,
        font="Aptos Display",
        bold=True,
    )
    add_text(
        slide,
        "Python Developer",
        Inches(0.85),
        Inches(1.72),
        Inches(5.0),
        Inches(0.55),
        size=20,
        color=TEAL,
        font="Aptos",
        bold=True,
    )
    add_divider(slide, 0.85, 2.35, 1.75)
    add_text(
        slide,
        "Building practical projects with clean code, useful libraries, and a strong learning mindset.",
        Inches(0.85),
        Inches(2.6),
        Inches(5.6),
        Inches(1.2),
        size=17,
        color=SLATE,
        font="Aptos",
    )

    add_chip(slide, "Python-first builder", 0.85, 4.0, 1.9, TEAL)
    add_chip(slide, "Libraries + projects", 2.95, 4.0, 2.1, DARK)
    add_chip(slide, "Open to internships", 5.3, 4.0, 2.0, AMBER)

    add_text(
        slide,
        "Focused on automation, backend basics, data handling, and continuous growth through real work.",
        Inches(0.85),
        Inches(4.72),
        Inches(6.3),
        Inches(0.9),
        size=14,
        color=SLATE,
        font="Aptos",
    )

    right = add_card(slide, Inches(8.0), Inches(0.65), Inches(4.45), Inches(6.2), DARK)
    right.fill.transparency = 0.0
    add_label(slide, "Profile Snapshot", 8.3, 0.95, w=2.05, h=0.38)

    add_card(slide, Inches(8.35), Inches(1.65), Inches(3.8), Inches(1.15), TEAL)
    add_text(
        slide,
        "Python Foundation",
        Inches(8.63),
        Inches(1.88),
        Inches(3.1),
        Inches(0.3),
        size=16,
        color=WHITE,
        font="Aptos Display",
        bold=True,
    )
    add_text(
        slide,
        "Clean syntax, scripting, logic, and beginner-friendly problem solving.",
        Inches(8.63),
        Inches(2.18),
        Inches(3.2),
        Inches(0.45),
        size=11,
        color=WHITE,
        font="Aptos",
    )

    add_card(slide, Inches(8.35), Inches(3.0), Inches(3.8), Inches(1.35), WHITE)
    add_text(
        slide,
        "Library Comfort",
        Inches(8.63),
        Inches(3.25),
        Inches(3.0),
        Inches(0.3),
        size=16,
        color=DARK,
        font="Aptos Display",
        bold=True,
    )
    add_text(
        slide,
        "NumPy, Pandas, Matplotlib,\nSeaborn, Requests, Flask",
        Inches(8.63),
        Inches(3.56),
        Inches(3.0),
        Inches(0.55),
        size=11,
        color=SLATE,
        font="Aptos",
    )

    add_card(slide, Inches(8.35), Inches(4.65), Inches(3.8), Inches(1.35), AMBER_SOFT)
    add_text(
        slide,
        "Career Direction",
        Inches(8.63),
        Inches(4.9),
        Inches(3.0),
        Inches(0.3),
        size=16,
        color=DARK,
        font="Aptos Display",
        bold=True,
    )
    add_text(
        slide,
        "Looking for internships, entry-level roles, and teams that value growth.",
        Inches(8.63),
        Inches(5.2),
        Inches(3.0),
        Inches(0.5),
        size=11,
        color=SLATE,
        font="Aptos",
    )

    add_text(
        slide,
        "github.com/Muneerali199",
        Inches(8.35),
        Inches(6.32),
        Inches(3.8),
        Inches(0.3),
        size=12,
        color=MUTED,
        font="Aptos",
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def slide_about():
    slide = prs.slides.add_slide(blank)
    add_bg(slide, DARK)
    add_blob(slide, Inches(-0.8), Inches(5.9), Inches(2.2), Inches(2.2), TEAL, transparency=0.15)
    add_blob(slide, Inches(11.7), Inches(-0.4), Inches(2.0), Inches(2.0), AMBER, transparency=0.12)

    add_label(slide, "About Me", 0.8, 0.65, w=1.6, h=0.38)
    add_text(
        slide,
        "A growing developer who learns by building.",
        Inches(0.8),
        Inches(1.2),
        Inches(6.4),
        Inches(0.9),
        size=27,
        color=WHITE,
        font="Aptos Display",
        bold=True,
    )
    add_text(
        slide,
        "I am at the start of my journey, but I care deeply about understanding the tools I use, writing readable code, and improving with every project.",
        Inches(0.8),
        Inches(2.15),
        Inches(5.9),
        Inches(1.1),
        size=15,
        color=MUTED,
        font="Aptos",
    )

    add_card(slide, Inches(0.8), Inches(3.55), Inches(2.0), Inches(2.0), WHITE)
    add_card(slide, Inches(3.0), Inches(3.55), Inches(2.0), Inches(2.0), TEAL_SOFT)
    add_card(slide, Inches(5.2), Inches(3.55), Inches(2.0), Inches(2.0), AMBER_SOFT)

    cards = [
        ("Builder\nMindset", "I prefer turning ideas into working projects so learning becomes practical."),
        ("Curious\nLearner", "I enjoy understanding libraries, not just copying syntax."),
        ("Steady\nGrowth", "Consistency matters to me more than quick shortcuts."),
    ]
    for idx, (title, body) in enumerate(cards):
        x = 1.05 + (idx * 2.2)
        fill_color = WHITE if idx == 0 else (TEAL_SOFT if idx == 1 else AMBER_SOFT)
        title_color = DARK if idx != 1 else TEAL
        add_text(
            slide,
            title,
            Inches(x),
            Inches(3.88),
            Inches(1.5),
            Inches(0.6),
            size=17,
            color=title_color,
            font="Aptos Display",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            body,
            Inches(x - 0.05),
            Inches(4.62),
            Inches(1.6),
            Inches(0.8),
            size=10.5,
            color=SLATE,
            font="Aptos",
            align=PP_ALIGN.CENTER,
        )

    quote = add_card(slide, Inches(8.0), Inches(1.05), Inches(4.5), Inches(5.6), WHITE)
    add_text(
        slide,
        "\"I enjoy learning useful tools and turning them into something real.\"",
        Inches(8.45),
        Inches(1.65),
        Inches(3.6),
        Inches(1.4),
        size=23,
        color=DARK,
        font="Aptos Display",
        bold=True,
    )
    add_divider(slide, 8.45, 3.1, 1.25, TEAL)
    add_paragraphs(
        slide,
        [
            "Python is my main language and my strongest foundation.",
            "I feel most motivated when I can connect code to a useful result.",
            "My goal is to keep growing into backend and project development with confidence.",
        ],
        Inches(8.45),
        Inches(3.45),
        Inches(3.55),
        Inches(2.2),
        size=13,
        color=SLATE,
        bullet=True,
    )


def slide_skills():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_blob(slide, Inches(11.8), Inches(0.2), Inches(1.8), Inches(1.8), TEAL_SOFT)
    add_blob(slide, Inches(-0.7), Inches(5.8), Inches(1.9), Inches(1.9), AMBER_SOFT)
    add_section_title(
        slide,
        "Skills",
        "Python + Library Stack",
        "A recruiter-friendly snapshot of what I can already work with and what I am actively learning next.",
    )

    columns = [
        (0.85, "Core Language And Workflow", TEAL, ["Python", "Git", "GitHub", "Markdown", "Jupyter", "VS Code"]),
        (4.45, "Libraries I Use", DARK, ["NumPy", "Pandas", "Matplotlib", "Seaborn", "Requests", "Flask"]),
        (8.05, "Currently Exploring", AMBER, ["SQL", "OOP", "APIs", "Problem Solving", "Backend Basics", "Project Structure"]),
    ]

    for x, heading, accent, items in columns:
        add_card(slide, Inches(x), Inches(2.0), Inches(3.1), Inches(3.9), WHITE)
        add_card(slide, Inches(x), Inches(2.0), Inches(3.1), Inches(0.12), accent, radius=False)
        add_text(
            slide,
            heading,
            Inches(x + 0.25),
            Inches(2.35),
            Inches(2.55),
            Inches(0.55),
            size=16,
            color=DARK,
            font="Aptos Display",
            bold=True,
        )
        chip_y = 3.1
        for idx, item in enumerate(items):
            add_chip(slide, item, x + 0.25, chip_y + idx * 0.45, 2.35, accent, WHITE)

    code = add_card(slide, Inches(0.85), Inches(6.2), Inches(10.3), Inches(0.75), DARK)
    add_text(
        slide,
        "focus = ['clean code', 'practical projects', 'steady improvement', 'strong fundamentals']",
        Inches(1.15),
        Inches(6.43),
        Inches(9.6),
        Inches(0.25),
        size=12,
        color=WHITE,
        font="Consolas",
    )


def slide_build_style():
    slide = prs.slides.add_slide(blank)
    add_bg(slide, TEAL_SOFT)
    add_blob(slide, Inches(11.6), Inches(5.7), Inches(2.3), Inches(2.3), WHITE, transparency=0.2)
    add_section_title(
        slide,
        "Workflow",
        "How I Like To Build",
        "My favorite learning loop is simple: understand the tool, use it in a project, then improve the result.",
    )

    steps = [
        ("01", "Learn The Tool", "I start by understanding the purpose of a library or concept before using it."),
        ("02", "Build Something Real", "I remember ideas better when I turn them into working scripts, notebooks, or small apps."),
        ("03", "Refine And Document", "I like making projects cleaner, easier to read, and better explained over time."),
    ]
    y_positions = [2.1, 3.55, 5.0]
    for (num, title, body), y in zip(steps, y_positions):
        add_card(slide, Inches(0.9), Inches(y), Inches(5.55), Inches(1.12), WHITE)
        add_card(slide, Inches(1.1), Inches(y + 0.18), Inches(0.62), Inches(0.62), TEAL)
        add_text(
            slide,
            num,
            Inches(1.23),
            Inches(y + 0.29),
            Inches(0.35),
            Inches(0.2),
            size=12,
            color=WHITE,
            font="Aptos",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            title,
            Inches(1.95),
            Inches(y + 0.18),
            Inches(2.2),
            Inches(0.3),
            size=16,
            color=DARK,
            font="Aptos Display",
            bold=True,
        )
        add_text(
            slide,
            body,
            Inches(1.95),
            Inches(y + 0.48),
            Inches(3.9),
            Inches(0.42),
            size=11.5,
            color=SLATE,
            font="Aptos",
        )

    add_card(slide, Inches(7.0), Inches(2.0), Inches(5.45), Inches(4.8), WHITE)
    add_text(
        slide,
        "What Excites Me Most",
        Inches(7.35),
        Inches(2.35),
        Inches(3.0),
        Inches(0.4),
        size=18,
        color=DARK,
        font="Aptos Display",
        bold=True,
    )

    boxes = [
        (7.35, 3.05, TEAL_SOFT, "Automation", "Saving time with scripts and repeatable logic."),
        (9.95, 3.05, AMBER_SOFT, "Data Tasks", "Cleaning, exploring, and presenting information clearly."),
        (7.35, 4.55, WHITE, "APIs", "Working with requests, responses, and useful integrations."),
        (9.95, 4.55, WHITE, "Backend Basics", "Learning how small services and apps fit together."),
    ]
    for x, y, fill, title, body in boxes:
        add_card(slide, Inches(x), Inches(y), Inches(2.25), Inches(1.15), fill)
        add_text(
            slide,
            title,
            Inches(x + 0.18),
            Inches(y + 0.18),
            Inches(1.8),
            Inches(0.25),
            size=13,
            color=TEAL if fill == TEAL_SOFT else DARK,
            font="Aptos Display",
            bold=True,
        )
        add_text(
            slide,
            body,
            Inches(x + 0.18),
            Inches(y + 0.45),
            Inches(1.85),
            Inches(0.45),
            size=10.2,
            color=SLATE,
            font="Aptos",
        )


def slide_roadmap():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_section_title(
        slide,
        "Roadmap",
        "Growth Plan For 2026",
        "A realistic path that turns current strengths into stronger project experience and team-ready confidence.",
    )

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.55), Inches(9.0), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = MUTED
    line.line.fill.background()

    milestones = [
        (1.2, TEAL, "Now", "Stronger Python projects\nand cleaner documentation"),
        (3.55, DARK, "Next", "Deeper SQL, OOP,\nAPIs, and backend basics"),
        (5.9, AMBER, "Then", "Larger collaborative\nor open-source work"),
        (8.25, TEAL, "Goal", "Internship or junior role\nwith meaningful learning"),
    ]
    for x, color, label, body in milestones:
        add_blob(slide, Inches(x), Inches(3.25), Inches(0.42), Inches(0.42), color)
        add_text(
            slide,
            label,
            Inches(x - 0.1),
            Inches(2.75),
            Inches(0.65),
            Inches(0.25),
            size=12,
            color=color,
            font="Aptos",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            body,
            Inches(x - 0.4),
            Inches(3.85),
            Inches(1.25),
            Inches(0.85),
            size=11.2,
            color=SLATE,
            font="Aptos",
            align=PP_ALIGN.CENTER,
        )

    add_card(slide, Inches(10.7), Inches(1.65), Inches(1.95), Inches(4.9), DARK)
    add_text(
        slide,
        "Best Fit",
        Inches(11.05),
        Inches(2.0),
        Inches(1.25),
        Inches(0.3),
        size=18,
        color=WHITE,
        font="Aptos Display",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_chip(slide, "Internships", 10.98, 2.7, 1.42, TEAL)
    add_chip(slide, "Junior Roles", 10.98, 3.22, 1.42, AMBER)
    add_chip(slide, "Mentorship", 10.98, 3.74, 1.42, TEAL)
    add_chip(slide, "Real Projects", 10.98, 4.26, 1.42, DARK)
    add_text(
        slide,
        "Looking for opportunities where curiosity, consistency, and growth are valued.",
        Inches(10.95),
        Inches(4.95),
        Inches(1.5),
        Inches(0.8),
        size=10.5,
        color=MUTED,
        font="Aptos",
        align=PP_ALIGN.CENTER,
    )


def slide_recruiter():
    slide = prs.slides.add_slide(blank)
    add_bg(slide, DARK)
    add_blob(slide, Inches(-0.8), Inches(-0.8), Inches(2.4), Inches(2.4), TEAL, transparency=0.15)
    add_blob(slide, Inches(11.4), Inches(5.6), Inches(2.5), Inches(2.5), AMBER, transparency=0.15)
    add_label(slide, "Recruiter Note", 0.8, 0.65, w=1.95, h=0.38)
    add_text(
        slide,
        "What You Can Expect From Me",
        Inches(0.8),
        Inches(1.2),
        Inches(6.0),
        Inches(0.75),
        size=27,
        color=WHITE,
        font="Aptos Display",
        bold=True,
    )

    add_paragraphs(
        slide,
        [
            "A genuine Python foundation with practical interest in useful libraries.",
            "A coachable mindset and willingness to improve quickly from feedback.",
            "Consistent effort, clear communication, and pride in readable work.",
            "Motivation to grow through real tasks, not only tutorials.",
        ],
        Inches(0.95),
        Inches(2.2),
        Inches(5.7),
        Inches(2.8),
        size=14,
        color=MUTED,
        bullet=True,
    )

    add_card(slide, Inches(7.35), Inches(1.35), Inches(4.55), Inches(4.95), WHITE)
    qualities = [
        ("Curiosity", "I like to understand the why behind the code."),
        ("Consistency", "I keep showing up and improving step by step."),
        ("Coachability", "I learn well in environments with feedback and direction."),
        ("Ownership", "I care about making projects cleaner and more useful over time."),
    ]
    y = 1.8
    for title, body in qualities:
        add_text(
            slide,
            title,
            Inches(7.75),
            Inches(y),
            Inches(1.7),
            Inches(0.25),
            size=15,
            color=TEAL,
            font="Aptos Display",
            bold=True,
        )
        add_text(
            slide,
            body,
            Inches(7.75),
            Inches(y + 0.24),
            Inches(3.5),
            Inches(0.42),
            size=10.8,
            color=SLATE,
            font="Aptos",
        )
        y += 0.92

    add_text(
        slide,
        "I may be early in my career, but I am serious about becoming excellent through disciplined work.",
        Inches(0.95),
        Inches(5.45),
        Inches(5.9),
        Inches(0.8),
        size=16,
        color=WHITE,
        font="Aptos",
        italic=True,
    )


def slide_close():
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_blob(slide, Inches(-0.9), Inches(5.7), Inches(2.3), Inches(2.3), TEAL_SOFT)
    add_blob(slide, Inches(11.8), Inches(-0.7), Inches(2.4), Inches(2.4), AMBER_SOFT)

    add_text(
        slide,
        "Open To Internships And Entry-Level Opportunities",
        Inches(1.15),
        Inches(1.15),
        Inches(11.0),
        Inches(0.7),
        size=28,
        color=DARK,
        font="Aptos Display",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "If you value growth, curiosity, and practical Python skills, I would be happy to connect.",
        Inches(2.05),
        Inches(2.0),
        Inches(9.2),
        Inches(0.5),
        size=15,
        color=SLATE,
        font="Aptos",
        align=PP_ALIGN.CENTER,
    )

    add_card(slide, Inches(3.85), Inches(3.05), Inches(5.65), Inches(1.35), DARK)
    add_text(
        slide,
        "github.com/Muneerali199",
        Inches(4.1),
        Inches(3.42),
        Inches(5.15),
        Inches(0.35),
        size=22,
        color=WHITE,
        font="Aptos Display",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_chip(slide, "Python", 4.15, 4.75, 1.0, TEAL)
    add_chip(slide, "Libraries", 5.35, 4.75, 1.1, DARK)
    add_chip(slide, "Projects", 6.65, 4.75, 1.0, AMBER)
    add_chip(slide, "Growth", 7.85, 4.75, 0.95, TEAL)

    add_text(
        slide,
        "Start small, stay consistent, and let the work speak for you.",
        Inches(2.2),
        Inches(6.0),
        Inches(8.9),
        Inches(0.4),
        size=17,
        color=TEAL,
        font="Aptos",
        italic=True,
        align=PP_ALIGN.CENTER,
    )


slide_cover()
slide_about()
slide_skills()
slide_build_style()
slide_roadmap()
slide_recruiter()
slide_close()

prs.save("Mishthi_Profile_Deck.pptx")
